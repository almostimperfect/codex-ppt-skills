#!/usr/bin/env python3
"""Package full-slide images into a PowerPoint deck."""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp"}


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Create an image-based PPTX from slide images.")
    parser.add_argument("--image-dir", required=True, help="Directory containing slide images.")
    parser.add_argument("--output", required=True, help="Output .pptx path.")
    parser.add_argument("--source-pptx", help="Optional source deck whose slide size should be reused.")
    parser.add_argument("--width", type=float, default=20.0, help="Slide width in inches when no source deck is given.")
    parser.add_argument("--height", type=float, default=11.25, help="Slide height in inches when no source deck is given.")
    return parser.parse_args()


def sorted_images(image_dir: Path) -> list[Path]:
    files = [p for p in image_dir.iterdir() if p.is_file() and p.suffix.lower() in IMAGE_EXTS]
    return sorted(files, key=lambda p: p.name)


def main() -> None:
    args = parse_args()
    image_dir = Path(args.image_dir)
    output = Path(args.output)
    images = sorted_images(image_dir)
    if not images:
        raise SystemExit(f"No slide images found in {image_dir}")

    if args.source_pptx:
        src = Presentation(args.source_pptx)
        prs = Presentation()
        prs.slide_width = src.slide_width
        prs.slide_height = src.slide_height
    else:
        prs = Presentation()
        prs.slide_width = Inches(args.width)
        prs.slide_height = Inches(args.height)

    blank = prs.slide_layouts[6]
    for image in images:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(image.resolve()), 0, 0, width=prs.slide_width, height=prs.slide_height)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"wrote {output} with {len(images)} slides")


if __name__ == "__main__":
    main()
