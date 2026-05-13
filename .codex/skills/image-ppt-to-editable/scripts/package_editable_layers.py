#!/usr/bin/env python3
"""Build a semi-editable PPTX from textless backgrounds and visual layout JSON.

Inputs:
- a directory of background images named slide-01.png, slide-02.png, ...
- a directory of layout JSON files named slide-01.json, slide-02.json, ...

The script places each background as a full-slide image, then adds editable
text boxes using pixel coordinates from the JSON layout.
"""

from __future__ import annotations

import argparse
import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Package textless slide backgrounds with editable text overlays.")
    parser.add_argument("--background-dir", required=True)
    parser.add_argument("--layout-dir", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument("--source-pptx", help="Optional source deck for slide size.")
    parser.add_argument("--width", type=float, default=13.333333, help="Slide width in inches when no source PPTX is provided.")
    parser.add_argument("--height", type=float, default=7.5, help="Slide height in inches when no source PPTX is provided.")
    parser.add_argument("--default-font", default="Aptos", help="Fallback font when a text item has no style.font_name.")
    return parser.parse_args()


def rgb(value: str) -> RGBColor:
    value = value.strip().lstrip("#")
    if len(value) != 6:
        return RGBColor(24, 48, 82)
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def align(value: str):
    return {
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "left": PP_ALIGN.LEFT,
    }.get((value or "left").lower(), PP_ALIGN.LEFT)


def main() -> None:
    args = parse_args()
    bg_dir = Path(args.background_dir)
    layout_dir = Path(args.layout_dir)
    output = Path(args.output)

    prs = Presentation()
    if args.source_pptx:
        src = Presentation(args.source_pptx)
        prs.slide_width = src.slide_width
        prs.slide_height = src.slide_height
    else:
        prs.slide_width = Inches(args.width)
        prs.slide_height = Inches(args.height)

    blank = prs.slide_layouts[6]
    slide_name = re.compile(r"^slide-\d+\.(png|jpg|jpeg|webp)$", re.IGNORECASE)
    backgrounds = sorted(p for p in bg_dir.iterdir() if p.is_file() and slide_name.match(p.name))
    if not backgrounds:
        raise SystemExit(f"No slide backgrounds found in {bg_dir}")

    for bg in backgrounds:
        layout_path = layout_dir / f"{bg.stem}.json"
        data = json.loads(layout_path.read_text(encoding="utf-8"))
        img_w = data["image_size"]["width"]
        img_h = data["image_size"]["height"]
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(bg.resolve()), 0, 0, width=prs.slide_width, height=prs.slide_height)

        for item in data.get("items", []):
            x, y, w, h = item["bbox"]
            style = item.get("style", {})
            box = slide.shapes.add_textbox(
                int(x / img_w * prs.slide_width),
                int(y / img_h * prs.slide_height),
                int(w / img_w * prs.slide_width),
                int(h / img_h * prs.slide_height),
            )
            tf = box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.margin_left = Pt(0)
            tf.margin_right = Pt(0)
            tf.margin_top = Pt(0)
            tf.margin_bottom = Pt(0)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = align(style.get("align", "left"))
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            run = p.add_run()
            run.text = item.get("text", "")
            run.font.name = style.get("font_name", args.default_font)
            run.font.size = Pt(float(style.get("font_size", 14)))
            run.font.bold = bool(style.get("bold", False))
            run.font.color.rgb = rgb(style.get("color", "#183052"))

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"wrote {output} with {len(backgrounds)} slides")


if __name__ == "__main__":
    main()
